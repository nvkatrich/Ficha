"""Turn attached specification files into a reviewable product list."""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Iterable

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from xlrd import open_workbook


class SpecificationError(RuntimeError):
    pass


HEADER_ALIASES = {
    "name": {"наименование", "название", "позиция", "товар", "услуга", "описание", "product", "item", "name"},
    "quantity": {"кол", "количество", "колво", "qty", "quantity"},
    "unit": {"ед", "единица", "едизм", "unit"},
    "price": {"цена", "цена за ед", "ценаед", "price", "unitprice"},
    "total": {"сумма", "стоимость", "итого", "total", "amount"},
    "comment": {"комментарий", "примечание", "comment", "note"},
}
SUPPORTED_EXTENSIONS = {".xlsx", ".xls", ".xlsm", ".csv", ".docx", ".pptx", ".pdf"}


def _norm(value: Any) -> str:
    return re.sub(r"[^a-zа-яё0-9]+", "", str(value or "").lower().replace("ё", "е"))


def _number(value: Any) -> str:
    return str(value or "").strip()


def _header_mapping(header: Iterable[Any]) -> dict[str, int]:
    mapping: dict[str, int] = {}
    for index, value in enumerate(header):
        token = _norm(value)
        for field, aliases in HEADER_ALIASES.items():
            if token in {_norm(alias) for alias in aliases} or any(_norm(alias) in token and len(_norm(alias)) > 3 for alias in aliases):
                mapping.setdefault(field, index)
                break
    return mapping


def _rows_from_matrix(matrix: list[list[Any]], source_name: str) -> dict[str, Any]:
    matrix = [[cell for cell in row] for row in matrix if any(str(cell or "").strip() for cell in row)]
    if not matrix:
        raise SpecificationError(f"В файле «{source_name}» не найдена таблица.")
    best_index, best_mapping = 0, {}
    for index, row in enumerate(matrix[:20]):
        candidate = _header_mapping(row)
        if len(candidate) > len(best_mapping):
            best_index, best_mapping = index, candidate
    if "name" not in best_mapping:
        raise SpecificationError(f"В файле «{source_name}» не найден столбец с наименованием.")
    products: list[dict[str, str]] = []
    warnings: list[str] = []
    for row_number, row in enumerate(matrix[best_index + 1:], best_index + 2):
        def get(field: str) -> str:
            index = best_mapping.get(field)
            return _number(row[index] if index is not None and index < len(row) else "")
        name = get("name")
        if not name or _norm(name) in {"итого", "всего", "total"}:
            continue
        item = {field: get(field) for field in ("name", "quantity", "unit", "price", "total", "comment")}
        item["source_row"] = str(row_number)
        if not item["quantity"]:
            item["quantity"] = "1"
            warnings.append(f"Строка {row_number}: количество не найдено, установлено 1.")
        if not item["price"] and not item["total"]:
            warnings.append(f"Строка {row_number}: нет цены или суммы, проверьте строку.")
        products.append(item)
    if not products:
        raise SpecificationError(f"В файле «{source_name}» не найдено ни одной позиции.")
    return {"source_name": source_name, "products": products, "warnings": warnings, "header": best_mapping}


def _matrix_from_xlsx(path: Path) -> list[list[Any]]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    for sheet in workbook.worksheets:
        rows = [list(row) for row in sheet.iter_rows(values_only=True)]
        if rows:
            return rows
    return []


def _matrix_from_xls(path: Path) -> list[list[Any]]:
    workbook = open_workbook(str(path), on_demand=True)
    for sheet in workbook.sheets():
        rows = [sheet.row_values(index) for index in range(sheet.nrows)]
        if rows:
            return rows
    return []


def _matrix_from_docx(path: Path) -> list[list[Any]]:
    document = Document(path)
    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        if rows:
            return rows
    raise SpecificationError(f"В файле «{path.name}» не найдена таблица.")


def _matrix_from_pptx(path: Path) -> list[list[Any]]:
    presentation = Presentation(path)
    candidates: list[list[list[Any]]] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                candidates.append([[cell.text for cell in row.cells] for row in shape.table.rows])
    if not candidates:
        raise SpecificationError(f"В файле «{path.name}» не найдена таблица PowerPoint.")
    return max(candidates, key=len)


def _matrix_from_pdf(path: Path) -> list[list[Any]]:
    rows: list[list[Any]] = []
    for page in PdfReader(str(path)).pages:
        text = page.extract_text() or ""
        for line in text.splitlines():
            cells = [cell.strip() for cell in re.split(r"\t+|\s{3,}", line) if cell.strip()]
            if len(cells) >= 2:
                rows.append(cells)
    return rows


def parse_specification(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise SpecificationError("Поддерживаются спецификации XLSX, XLS, CSV, DOCX, PPTX и текстовые PDF.")
    try:
        if suffix in {".xlsx", ".xlsm"}:
            result = _rows_from_matrix(_matrix_from_xlsx(path), path.name)
        elif suffix == ".xls":
            result = _rows_from_matrix(_matrix_from_xls(path), path.name)
        elif suffix == ".docx":
            result = _rows_from_matrix(_matrix_from_docx(path), path.name)
        elif suffix == ".pptx":
            result = _rows_from_matrix(_matrix_from_pptx(path), path.name)
        elif suffix == ".csv":
            import csv
            with path.open("r", encoding="utf-8-sig", newline="") as file:
                result = _rows_from_matrix([row for row in csv.reader(file)], path.name)
        else:
            result = _rows_from_matrix(_matrix_from_pdf(path), path.name)
    except SpecificationError:
        raise
    except Exception as error:
        raise SpecificationError(f"Не удалось разобрать спецификацию «{path.name}».") from error
    result["extension"] = suffix
    return result
