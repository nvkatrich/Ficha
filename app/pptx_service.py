"""Native PowerPoint template generation.

The template is edited as a PowerPoint package: text boxes remain text boxes and
product/budget grids remain native PowerPoint tables rather than raster images.
"""
from __future__ import annotations

import copy
import re
import uuid
from pathlib import Path
from typing import Any

from pptx import Presentation

from .config import DOCUMENT_DIR
from .document_service import as_text, lookup_dotted
from .template_store import copy_template_for_edit

TOKEN_PATTERN = re.compile(r"{{\s*([A-Za-zА-Яа-яЁё0-9_.]+)\s*}}")


class PowerPointError(RuntimeError):
    pass


def _replace_text_frame(text_frame: Any, context: dict[str, Any]) -> bool:
    runs = [run for paragraph in text_frame.paragraphs for run in paragraph.runs]
    full_text = "".join(run.text for run in runs)
    if not TOKEN_PATTERN.search(full_text):
        return False
    replaced = TOKEN_PATTERN.sub(lambda match: lookup_dotted(context, match.group(1)), full_text)
    if runs:
        runs[0].text = replaced
        for run in runs[1:]:
            run.text = ""
    else:
        text_frame.text = replaced
    return True


def _set_shape_text(shape: Any, text: str) -> None:
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    runs = [run for paragraph in frame.paragraphs for run in paragraph.runs]
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        frame.text = text


def _set_cell_text(cell: Any, text: Any) -> None:
    if hasattr(cell, "has_text_frame"):
        _set_shape_text(cell, as_text(text))
    else:
        _set_text_frame(cell.text_frame, as_text(text))


def _set_text_frame(frame: Any, text: str) -> None:
    runs = [run for paragraph in frame.paragraphs for run in paragraph.runs]
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        frame.text = text


def _iter_shapes(shapes: Any):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == 6:  # group shape
            yield from _iter_shapes(shape.shapes)


def _add_cloned_row(table: Any, before_row: int) -> Any:
    """Clone the existing last data row so formatting survives row expansion."""
    rows_xml = table._tbl.tr_lst
    template_index = max(1, before_row - 1)
    new_row = copy.deepcopy(rows_xml[template_index])
    table._tbl.insert(before_row, new_row)
    return table.rows[before_row]


def _remove_row(table: Any, index: int) -> None:
    table._tbl.remove(table._tbl.tr_lst[index])


def _fill_product_table(shape: Any, products: list[dict[str, Any]], total: str) -> None:
    table = shape.table
    if len(table.rows) < 2:
        return
    total_row_index = len(table.rows) - 1
    desired_data_rows = max(1, len(products))
    while len(table.rows) < desired_data_rows + 2:
        _add_cloned_row(table, len(table.rows) - 1)
    while len(table.rows) > desired_data_rows + 2:
        _remove_row(table, len(table.rows) - 2)
    total_row_index = len(table.rows) - 1
    for index, product in enumerate(products):
        row = table.rows[index + 1]
        values = [product.get("name", ""), product.get("quantity", ""), product.get("price", ""), product.get("total", "")]
        for cell, value in zip(row.cells, values):
            _set_cell_text(cell, value)
    if not products:
        for cell in table.rows[1].cells:
            _set_cell_text(cell, "")
    total_row = table.rows[total_row_index]
    if len(total_row.cells) >= 1:
        _set_cell_text(total_row.cells[0], "Итого оборудование и ПО")
    if len(total_row.cells) >= 4:
        _set_cell_text(total_row.cells[len(total_row.cells) - 1], total)


def _fill_budget_table(shape: Any, context: dict[str, Any]) -> None:
    table = shape.table
    if len(table.rows) < 4 or len(table.columns) < 2:
        return
    _set_cell_text(table.cell(1, 1), context.get("products_total", ""))
    # Keep the template's manually entered installation row intact. The project
    # total is the CRM opportunity and remains fully editable in PowerPoint.
    _set_cell_text(table.cell(len(table.rows) - 1, 1), context.get("deal", {}).get("amount", ""))


def _apply_named_bindings(shape: Any, context: dict[str, Any]) -> None:
    name = shape.name.lower()
    bindings = {
        "cover-title": "deal.title",
        "cover-customer": "client.company",
        "total-number": "deal.amount",
        "contact-name": "manager.name",
        "contact-company": "seller.name",
    }
    for shape_name, key in bindings.items():
        context_value = lookup_dotted(context, key)
        if shape_name in name and context_value:
            _set_shape_text(shape, context_value)
    if "contact-details" in name:
        phone = lookup_dotted(context, "manager.phone")
        email = lookup_dotted(context, "manager.email")
        if phone or email:
            _set_shape_text(shape, " · ".join(filter(None, (phone, email))))
    if "footer" in name:
        seller = lookup_dotted(context, "seller.name")
        if seller:
            _set_shape_text(shape, f"{seller} · Коммерческое предложение")


def generate_pptx(template: dict[str, Any], context: dict[str, Any], output_basename: str) -> Path:
    safe_name = re.sub(r"[^A-Za-zА-Яа-яЁё0-9_.-]+", "_", output_basename).strip("._") or "commercial_proposal"
    output_path = DOCUMENT_DIR / f"{safe_name}_{uuid.uuid4().hex[:8]}.pptx"
    copy_template_for_edit(template, output_path)
    try:
        presentation = Presentation(output_path)
        for slide in presentation.slides:
            for shape in _iter_shapes(slide.shapes):
                if getattr(shape, "has_table", False):
                    name = shape.name.lower()
                    if "hardware" in name or "product" in name or "комплект" in name:
                        _fill_product_table(shape, context.get("products", []), context.get("products_total", ""))
                    elif "budget" in name or "стоим" in name or "summary" in name:
                        _fill_budget_table(shape, context)
                    for row in shape.table.rows:
                        for cell in row.cells:
                            _replace_text_frame(cell.text_frame, context)
                elif getattr(shape, "has_text_frame", False):
                    _apply_named_bindings(shape, context)
                    _replace_text_frame(shape.text_frame, context)
        presentation.save(output_path)
    except Exception as error:
        output_path.unlink(missing_ok=True)
        raise PowerPointError("Не удалось собрать редактируемую презентацию из этого шаблона.") from error
    return output_path
