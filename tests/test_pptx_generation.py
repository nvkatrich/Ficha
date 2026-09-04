from pathlib import Path

from pptx import Presentation

from app import document_service
from app import pptx_service
from app import template_store


SOURCE = Path("/home/ubuntu/upload/pasted_file_6pDAqu_Коммерческое_предложение_Корпорация_Маяк_Екатеринбург.pptx")


def test_provided_pptx_stays_native_and_tables_are_editable(tmp_path, monkeypatch):
    templates = tmp_path / "templates"
    documents = tmp_path / "documents"
    templates.mkdir()
    documents.mkdir()
    target = templates / "source.pptx"
    target.write_bytes(SOURCE.read_bytes())
    monkeypatch.setattr(template_store, "TEMPLATE_DIR", templates)
    monkeypatch.setattr(pptx_service, "DOCUMENT_DIR", documents)

    context = document_service.build_context(
        {
            "deal": {"ID": "410", "TITLE": "Автоматизированная парковочная система", "OPPORTUNITY": "2095500", "CURRENCY_ID": "RUB"},
            "company": {"TITLE": "Корпорация Маяк Екатеринбург"},
            "contact": {},
            "manager": {"NAME": "Анна", "LAST_NAME": "Иванова", "EMAIL": "a@example.ru", "WORK_PHONE": "+7 900 000-00-00"},
            "products": [
                {"PRODUCT_NAME": "Въездная стойка", "QUANTITY": "1", "PRICE": "365000", "PRICE_ACCOUNT": "365000", "CURRENCY": "RUB"},
                {"PRODUCT_NAME": "Камера распознавания ГРЗ", "QUANTITY": "2", "PRICE": "36500", "PRICE_ACCOUNT": "73000", "CURRENCY": "RUB"},
            ],
        },
        {"document_prefix": "КП", "company_name": "ООО Новые Технологии"},
    )

    output = pptx_service.generate_pptx({"filename": "source.pptx"}, context, "Маяк")
    presentation = Presentation(output)
    assert len(presentation.slides) == 5
    hardware = next(shape for shape in presentation.slides[3].shapes if shape.name == "hardware-table")
    assert hardware.has_table
    assert len(hardware.table.rows) == 4  # header + 2 native data rows + total
    assert hardware.table.cell(1, 0).text == "Въездная стойка"
    assert hardware.table.cell(2, 0).text == "Камера распознавания ГРЗ"
    assert hardware.table.cell(3, 3).text == "438 000,00 RUB"
    budget = next(shape for shape in presentation.slides[4].shapes if shape.name == "budget-table")
    assert budget.table.cell(1, 1).text == "438 000,00 RUB"
    assert presentation.slides[0].shapes[4].text == "Автоматизированная парковочная система"
    all_text = "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "{{" not in all_text
